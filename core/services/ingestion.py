"""
Document Ingestion Pipeline Service.

Implements Schema #2 from ARCHITECTURE.md:
    Document → Extractor → Raw Text → Cleaner → Chunker → Database

This service extracts text and metadata from uploaded documents (PDF, DOCX, PPTX),
normalizes the text for search, and stores chunked content
for retrieval-augmented generation (RAG).
"""

import logging
import os
import re
import tempfile
from typing import List, Tuple

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation

from core.models import Document, Page, Chunk


logger = logging.getLogger(__name__)

# Chunk configuration
CHUNK_SIZE = 1000  # Target chunk size in characters
CHUNK_OVERLAP = 100  # Overlap between chunks for context continuity


class DocumentProcessingError(Exception):
    """Raised when document processing fails (corrupted file, unsupported format, etc.)."""
    pass


# Backward compatibility alias
PDFProcessingError = DocumentProcessingError


def clean_text(raw_text: str) -> str:
    """
    Clean extracted text by removing excessive whitespace.
    
    Args:
        raw_text: Text extracted directly from PDF.
        
    Returns:
        Cleaned text with normalized whitespace.
    """
    # Replace multiple whitespace/newlines with single space
    cleaned = re.sub(r'\s+', ' ', raw_text)
    # Strip leading/trailing whitespace
    return cleaned.strip()


def normalize_text(text: str) -> str:
    """
    Normalize text for case-insensitive search.
    
    Args:
        text: Cleaned text.
        
    Returns:
        Lowercase text for search matching.
    """
    return text.lower()


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """
    Split text into overlapping chunks of approximately `chunk_size` characters.
    
    Chunks are created with overlap to preserve context at boundaries.
    The algorithm tries to break at sentence boundaries when possible.
    
    Args:
        text: The text to chunk.
        chunk_size: Target size for each chunk (default: 1000 chars).
        overlap: Number of characters to overlap between chunks (default: 100).
        
    Returns:
        List of text chunks.
    """
    if not text or len(text) <= chunk_size:
        return [text] if text else []
    
    chunks = []
    start = 0
    text_length = len(text)
    
    while start < text_length:
        # Calculate end position
        end = start + chunk_size
        
        if end >= text_length:
            # Last chunk - take everything remaining
            chunks.append(text[start:])
            break
        
        # Try to find a good breaking point (sentence boundary)
        # Look for period, question mark, or exclamation followed by space
        search_start = max(start + chunk_size - 200, start)  # Search in last 200 chars
        search_region = text[search_start:end + 50]  # Look slightly beyond target
        
        # Find the last sentence boundary in the search region
        best_break = None
        for match in re.finditer(r'[.!?]\s', search_region):
            best_break = search_start + match.end()
        
        if best_break and best_break > start + chunk_size // 2:
            # Use sentence boundary if it's not too early in the chunk
            end = best_break
        
        chunks.append(text[start:end].strip())
        
        # Move start position, accounting for overlap
        start = end - overlap
    
    return chunks


def extract_text_from_docx(file_path: str) -> List[Tuple[str, float, float, int]]:
    """
    Extract text from a DOCX file.
    
    Args:
        file_path: Path to the .docx file.
        
    Returns:
        List of tuples: (text, width, height, page_number)
        For DOCX, width/height are set to standard A4 dimensions (612x792 points).
        Rotation is always 0.
        
    Raises:
        DocumentProcessingError: If the file cannot be opened or is corrupted.
    """
    try:
        doc = DocxDocument(file_path)
        
        # Extract all paragraphs
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text)
        
        # Combine all text
        full_text = '\n'.join(paragraphs)
        
        # For DOCX, we treat the entire document as one "page"
        # Use standard A4 dimensions (8.5 x 11 inches = 612 x 792 points)
        if full_text:
            return [(full_text, 612.0, 792.0, 1)]
        return []
        
    except Exception as e:
        error_msg = f"Failed to extract text from DOCX: {str(e)}"
        logger.error(error_msg)
        raise DocumentProcessingError(error_msg) from e


def _extract_text_from_shape(shape) -> List[str]:
    """
    Recursively extract text from a shape, handling nested groups.
    
    Args:
        shape: A PowerPoint shape object.
        
    Returns:
        List of text strings extracted from the shape.
    """
    text_parts = []
    
    try:
        # Handle text_frame (most common case)
        if hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if para_text:
                    text_parts.append(para_text)
        
        # Handle direct text attribute (fallback)
        elif hasattr(shape, "text") and shape.text:
            text = shape.text.strip()
            if text:
                text_parts.append(text)
        
        # Handle tables
        if hasattr(shape, "has_table") and shape.has_table:
            try:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            text_parts.append(cell_text)
            except Exception as e:
                logger.warning(f"Error extracting text from table: {e}")
        
        # Handle groups (nested shapes)
        if hasattr(shape, "shapes"):
            try:
                for sub_shape in shape.shapes:
                    text_parts.extend(_extract_text_from_shape(sub_shape))
            except Exception as e:
                logger.warning(f"Error extracting text from group: {e}")
                
    except Exception as e:
        logger.warning(f"Error extracting text from shape: {e}")
    
    return text_parts


def extract_text_from_pptx(file_path: str) -> List[Tuple[str, float, float, int]]:
    """
    Extract text from a PPTX file with robust error handling.
    
    Args:
        file_path: Path to the .pptx file.
        
    Returns:
        List of tuples: (text, width, height, slide_number)
        For PPTX, width/height are set to standard presentation dimensions (720x540 points).
        Rotation is always 0.
        
    Raises:
        DocumentProcessingError: If the file cannot be opened or is corrupted.
    """
    logger.info(f"Starting PPTX extraction from: {file_path}")
    print(f"[PPTX] Starting extraction from: {file_path}")  # Debug print
    
    try:
        # Open presentation
        logger.debug("Opening PPTX file...")
        print("[PPTX] Opening file...")  # Debug print
        prs = Presentation(file_path)
        total_slides = len(prs.slides)
        logger.info(f"Found {total_slides} slides in PPTX file")
        print(f"[PPTX] Found {total_slides} slides")  # Debug print
        
        slides_data = []
        
        # Iterate through all slides with error handling per slide
        for slide_num, slide in enumerate(prs.slides, start=1):
            try:
                logger.debug(f"Processing slide {slide_num}/{total_slides}...")
                print(f"[PPTX] Processing slide {slide_num}/{total_slides}...")  # Debug print
                slide_text_parts = []
                
                # Extract text from all shapes on the slide
                for shape_idx, shape in enumerate(slide.shapes):
                    try:
                        shape_texts = _extract_text_from_shape(shape)
                        slide_text_parts.extend(shape_texts)
                    except Exception as e:
                        # Log but continue - one broken shape shouldn't stop the whole slide
                        logger.warning(
                            f"Error extracting text from shape {shape_idx} on slide {slide_num}: {e}"
                        )
                        continue
                
                # Combine all text from this slide
                slide_text = '\n'.join(slide_text_parts)
                
                if slide_text:
                    # Use standard presentation dimensions (10 x 7.5 inches = 720 x 540 points)
                    slides_data.append((slide_text, 720.0, 540.0, slide_num))
                    logger.debug(f"Slide {slide_num}: Extracted {len(slide_text)} characters")
                else:
                    logger.debug(f"Slide {slide_num}: No text found")
                    
            except Exception as e:
                # Log slide-level errors but continue processing other slides
                logger.warning(f"Error processing slide {slide_num}: {e}")
                continue
        
        logger.info(f"PPTX extraction complete: {len(slides_data)} slides with text")
        print(f"[PPTX] Extraction complete: {len(slides_data)} slides with text")  # Debug print
        return slides_data
        
    except Exception as e:
        error_msg = f"Failed to extract text from PPTX: {str(e)}"
        logger.error(error_msg, exc_info=True)
        raise DocumentProcessingError(error_msg) from e


def process_document(doc_id: int) -> dict:
    """
    Process a document: extract text, metadata, and create chunks.
    
    Supports PDF, DOCX, and PPTX file formats.
    This is the main entry point for the ingestion pipeline.
    
    Args:
        doc_id: Primary key of the Document to process.
        
    Returns:
        Dictionary with processing statistics:
        {
            'success': bool,
            'pages_processed': int,
            'chunks_created': int,
            'error': str | None
        }
        
    Raises:
        DocumentProcessingError: If the document cannot be opened or is corrupted.
        Document.DoesNotExist: If no document exists with the given ID.
    """
    result = {
        'success': False,
        'pages_processed': 0,
        'chunks_created': 0,
        'error': None
    }
    
    # Step 1: Retrieve the Document
    try:
        document = Document.objects.get(pk=doc_id)
    except Document.DoesNotExist:
        logger.error(f"Document with ID {doc_id} not found")
        raise
    
    # Get file path - handle both saved files and in-memory uploads
    temp_file_created = False
    try:
        file_path = document.file.path
        logger.debug(f"Using file path: {file_path}")
    except (ValueError, AttributeError):
        # File might not be saved to disk yet, use temporary file
        logger.warning("File not saved to disk, using temporary file")
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(document.file.name)[1]) as tmp_file:
            for chunk in document.file.chunks():
                tmp_file.write(chunk)
            file_path = tmp_file.name
            temp_file_created = True
        logger.debug(f"Created temporary file: {file_path}")
    
    file_extension = os.path.splitext(file_path)[1].lower()
    
    logger.info(f"Starting document processing: {document.title} (ID: {doc_id}, Type: {file_extension})")
    
    # Step 2: Extract text based on file type
    pages_data = []  # List of (text, width, height, page_number)
    temp_file_created = False
    
    try:
        # Verify file exists and is readable
        if not os.path.exists(file_path):
            raise DocumentProcessingError(f"File not found: {file_path}")
        
        if not os.access(file_path, os.R_OK):
            raise DocumentProcessingError(f"File not readable: {file_path}")
        
        logger.debug(f"File size: {os.path.getsize(file_path)} bytes")
        
        if file_extension == '.pdf':
            # Process PDF using PyMuPDF
            logger.debug("Processing as PDF...")
            pdf_doc = fitz.open(file_path)
            try:
                for page_num in range(len(pdf_doc)):
                    pdf_page = pdf_doc[page_num]
                    raw_text = pdf_page.get_text()
                    rect = pdf_page.rect
                    pages_data.append((raw_text, rect.width, rect.height, page_num + 1))
            finally:
                pdf_doc.close()
                
        elif file_extension == '.docx':
            # Process DOCX
            logger.debug("Processing as DOCX...")
            pages_data = extract_text_from_docx(file_path)
            
        elif file_extension == '.pptx':
            # Process PPTX
            logger.debug("Processing as PPTX...")
            pages_data = extract_text_from_pptx(file_path)
            
        else:
            error_msg = f"Unsupported file format: {file_extension}. Supported: .pdf, .docx, .pptx"
            logger.error(error_msg)
            result['error'] = error_msg
            raise DocumentProcessingError(error_msg)
        
        # Clean up temporary file if created
        if temp_file_created and os.path.exists(file_path):
            try:
                os.unlink(file_path)
                logger.debug(f"Cleaned up temporary file: {file_path}")
            except Exception as e:
                logger.warning(f"Failed to delete temporary file: {e}")
        
        # Step 3: Process each page/slide
        for text, width, height, page_num in pages_data:
            try:
                cleaned_text = clean_text(text)
                normalized_text = normalize_text(cleaned_text)
                
                # Create Page record (1-indexed page numbers)
                page = Page.objects.create(
                    document=document,
                    page_number=page_num,
                    text_raw=cleaned_text,
                    text_norm=normalized_text,
                    width=width,
                    height=height,
                    rotation=0  # DOCX and PPTX don't have rotation
                )
                
                result['pages_processed'] += 1
                logger.debug(f"Processed page/slide {page_num}: {len(cleaned_text)} chars")
                
                # Step 4: Chunk the page text
                if cleaned_text:
                    text_chunks = chunk_text(cleaned_text)
                    
                    for chunk_index, chunk_text_content in enumerate(text_chunks):
                        Chunk.objects.create(
                            page=page,
                            chunk_index=chunk_index,
                            chunk_text=chunk_text_content
                        )
                        result['chunks_created'] += 1
                
            except Exception as e:
                # Log page-level errors but continue processing other pages
                logger.warning(f"Error processing page {page_num}: {str(e)}")
                continue
        
        result['success'] = True
        logger.info(
            f"Document processing complete: {result['pages_processed']} pages/slides, "
            f"{result['chunks_created']} chunks"
        )
        
    except DocumentProcessingError:
        raise
    except Exception as e:
        error_msg = f"Error during document processing: {str(e)}"
        logger.error(error_msg)
        result['error'] = error_msg
        raise DocumentProcessingError(error_msg) from e
    
    return result


# Backward compatibility alias
def process_pdf_document(doc_id: int) -> dict:
    """
    Legacy function name for backward compatibility.
    
    Now calls process_document() which supports multiple formats.
    """
    return process_document(doc_id)

